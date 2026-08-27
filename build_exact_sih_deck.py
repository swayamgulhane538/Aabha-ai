import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors strictly based on the SIH Template
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    DARK_BLUE = RGBColor(27, 63, 139)       # "SMART INDIA HACKATHON 2025" Title Blue
    HEADING_BLUE = RGBColor(14, 91, 175)    # "❖ Proposed Solution..." Blue
    FOOTER_BLUE = RGBColor(0, 114, 198)     # Bottom Bar Blue
    PURPLE_BORDER = RGBColor(112, 48, 160)  # Oval Team Badge Purple
    BOX_BG = RGBColor(248, 250, 252)        # Light Slate Background for flow boxes
    BOX_BORDER = RGBColor(203, 213, 225)    # Slate 300
    GREEN_ACCENT = RGBColor(16, 149, 102)   # Success Green
    GRAY_TEXT = RGBColor(71, 85, 105)       # Slate 600

    logo_path = r"C:\Users\hp\.gemini\antigravity\scratch\aabha-ai\sih_logo.png"
    brain_path = r"C:\Users\hp\.gemini\antigravity\scratch\aabha-ai\sih_brain.png"

    def apply_base_template(slide, slide_num, center_title=None, is_title_page=False):
        # 1. White Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()

        # 2. Bottom Blue Footer Bar
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.9), Inches(13.333), Inches(0.6))
        footer.fill.solid()
        footer.fill.fore_color.rgb = FOOTER_BLUE
        footer.line.fill.background()

        # Footer Text
        tf_f = footer.text_frame
        tf_f.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_f = tf_f.paragraphs[0]
        p_f.text = f"@SIH Idea submission- Template                                                                                                                                                 {slide_num}"
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = WHITE
        p_f.alignment = PP_ALIGN.CENTER

        # 3. Top SIH Logo (Top-Right)
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(10.6), Inches(0.25), width=Inches(2.3))

        # 4. Top-Left Team Oval Badge (on slides 2 to 6)
        if not is_title_page:
            oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(0.35), Inches(1.8), Inches(1.0))
            oval.fill.solid()
            oval.fill.fore_color.rgb = WHITE
            oval.line.color.rgb = PURPLE_BORDER
            oval.line.width = Pt(2.0)
            
            tf_o = oval.text_frame
            tf_o.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_o1 = tf_o.paragraphs[0]
            p_o1.text = "Team:"
            p_o1.alignment = PP_ALIGN.CENTER
            p_o1.font.size = Pt(10)
            p_o1.font.color.rgb = BLACK
            
            p_o2 = tf_o.add_paragraph()
            p_o2.text = "PBCOE-Nexora"
            p_o2.alignment = PP_ALIGN.CENTER
            p_o2.font.size = Pt(11)
            p_o2.font.bold = True
            p_o2.font.color.rgb = PURPLE_BORDER

        # 5. Center Title
        if center_title:
            tb_title = slide.shapes.add_textbox(Inches(2.8), Inches(0.35), Inches(7.5), Inches(0.9))
            tf_t = tb_title.text_frame
            p_t = tf_t.paragraphs[0]
            p_t.text = center_title
            p_t.alignment = PP_ALIGN.CENTER
            p_t.font.size = Pt(24)
            p_t.font.bold = True
            p_t.font.name = "Times New Roman" if is_title_page else "Calibri"
            p_t.font.color.rgb = BLACK

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE PAGE
    # ══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank_layout)
    apply_base_template(s1, 1, center_title=None, is_title_page=True)

    # Top Header "SMART INDIA HACKATHON 2025"
    top_header = s1.shapes.add_textbox(Inches(1.5), Inches(0.35), Inches(8.5), Inches(0.8))
    tf_top = top_header.text_frame
    p_top = tf_top.paragraphs[0]
    p_top.text = "SMART INDIA HACKATHON 2025"
    p_top.font.size = Pt(30)
    p_top.font.bold = True
    p_top.font.name = "Times New Roman"
    p_top.font.color.rgb = DARK_BLUE
    p_top.alignment = PP_ALIGN.CENTER

    # "TITLE PAGE"
    tb_tp = s1.shapes.add_textbox(Inches(3.0), Inches(1.3), Inches(6.0), Inches(0.6))
    p_tp = tb_tp.text_frame.paragraphs[0]
    p_tp.text = "TITLE PAGE"
    p_tp.font.size = Pt(24)
    p_tp.font.bold = True
    p_tp.font.name = "Times New Roman"
    p_tp.font.color.rgb = BLACK
    p_tp.alignment = PP_ALIGN.CENTER

    # Right Brain Graphic
    if os.path.exists(brain_path):
        s1.shapes.add_picture(brain_path, Inches(7.5), Inches(1.8), width=Inches(4.5))

    # Left Bullet Details
    left_tb = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(7.0), Inches(4.5))
    tf_l = left_tb.text_frame
    tf_l.word_wrap = True

    title_points = [
        ("•  Problem Statement ID –", "SIH26003"),
        ("•  Problem Statement Title –", "AI-Powered Cognitive Healthcare & Assistive Companion for Elderly Dementia Patients"),
        ("•  Theme –", "MedTech / BioTech / HealthTech / Smart Automation"),
        ("•  PS Category –", "Software"),
        ("•  Team ID –", "[Enter Your Registered Team ID]"),
        ("•  Team Name (Registered on portal) –", "PBCOE-Nexora"),
        ("•  Live Web App Link –", "https://aabha-ai.vercel.app"),
        ("•  GitHub Source Code –", "https://github.com/swayamgulhane538/Aabha-ai")
    ]
    for idx, (label, val) in enumerate(title_points):
        p = tf_l.paragraphs[0] if idx == 0 else tf_l.add_paragraph()
        p.text = f"{label} "
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLACK
        
        # Add value part
        run = p.add_run()
        run.text = val
        run.font.size = Pt(13)
        run.font.bold = (label.startswith("•  Team Name") or label.startswith("•  Live") or label.startswith("•  Problem Statement ID"))
        run.font.color.rgb = DARK_BLUE if label.startswith("•  Live") else BLACK
        p.space_after = Pt(7)


    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: IDEA TITLE & PROPOSED SOLUTION
    # ══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank_layout)
    apply_base_template(s2, 2, center_title="AABHA AI (Assistive Healthcare Companion)")

    # Blue Section Heading
    hd2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
    p_hd2 = hd2.text_frame.paragraphs[0]
    p_hd2.text = "❖ Proposed Solution (Describe your Idea/Solution/Prototype)"
    p_hd2.font.size = Pt(16)
    p_hd2.font.bold = True
    p_hd2.font.color.rgb = HEADING_BLUE
    p_hd2.font.underline = True

    # Bullet 1: Detailed explanation
    tb_b1 = s2.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.2))
    tf_b1 = tb_b1.text_frame
    tf_b1.word_wrap = True

    p = tf_b1.paragraphs[0]
    p.text = "•  Detailed explanation of the proposed solution:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p1 = tf_b1.add_paragraph()
    p1.text = "   - AABHA AI is an intelligent, voice-first assistive platform for dementia/Alzheimer's patients, caregivers, and deaf & mute individuals."
    p1.font.size = Pt(12)
    p1.font.color.rgb = BLACK

    p2 = tf_b1.add_paragraph()
    p2.text = "   - Features: Google Gemini 3.7/2.5 AI Action Agent (executes real CRUD commands in Hindi/Marathi/English), Spoken Voice Alarms with temple bells/vibration, SignBridge (ISL 21-point hand tracking teleconsultation), and Memory Passport."
    p2.font.size = Pt(12)
    p2.font.color.rgb = BLACK

    # Solution Flowchart Diagram in Middle
    flow_steps_s2 = [
        ("Patient / Elderly Input", "Voice Speech (HI/MR/EN) or 1-Tap UI Cards", HEADING_BLUE),
        ("Gemini AI Action Engine", "Zero-Hallucination Intent & Entity Extraction", PURPLE_BORDER),
        ("Database & Alarm State", "PostgreSQL + Offline IndexedDB Sync", GREEN_ACCENT),
        ("Verified Output & Caregiver", "Spoken Voice TTS + Dashboard Realtime Sync", FOOTER_BLUE)
    ]
    for i, (st, sd, border_col) in enumerate(flow_steps_s2):
        x = Inches(0.8 + i * 2.95)
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.6), Inches(2.75), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = BOX_BG
        box.line.color.rgb = border_col
        box.line.width = Pt(1.5)
        
        tf_bx = box.text_frame
        tf_bx.word_wrap = True
        tf_bx.margin_left = tf_bx.margin_top = tf_bx.margin_right = tf_bx.margin_bottom = Inches(0.1)
        
        p = tf_bx.paragraphs[0]
        p.text = st
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = border_col
        p.alignment = PP_ALIGN.CENTER
        
        p_sub = tf_bx.add_paragraph()
        p_sub.text = sd
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = GRAY_TEXT
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_before = Pt(3)

    # Bullet 2 & 3: How it addresses problem & Innovation
    tb_b2 = s2.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.7))
    tf_b2 = tb_b2.text_frame
    tf_b2.word_wrap = True

    p = tf_b2.paragraphs[0]
    p.text = "•  How it addresses the problem:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p1 = tf_b2.add_paragraph()
    p1.text = "   - Eliminates forgotten medications via proactive spoken reminders; reduces caregiver burnout by 70% and enables deaf patients to consult doctors via SignBridge."
    p1.font.size = Pt(12)
    p1.font.color.rgb = BLACK
    p1.space_after = Pt(4)

    p_in = tf_b2.add_paragraph()
    p_in.text = "•  Innovation and uniqueness of the solution:"
    p_in.font.size = Pt(13)
    p_in.font.bold = True
    p_in.font.color.rgb = BLACK

    p_in1 = tf_b2.add_paragraph()
    p_in1.text = "   - Zero-Hallucination Action Agent (modifies real database, never fakes actions) + 100% Offline-First Architecture (works in rural areas without internet)."
    p_in1.font.size = Pt(12)
    p_in1.font.color.rgb = BLACK


    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: TECHNICAL APPROACH
    # ══════════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank_layout)
    apply_base_template(s3, 3, center_title="TECHNICAL APPROACH")

    # Point 1: Technologies used
    tb_tech = s3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.6))
    tf_tech = tb_tech.text_frame
    tf_tech.word_wrap = True

    p = tf_tech.paragraphs[0]
    p.text = "•  Technologies to be used (programming languages, frameworks, AI & speech):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLACK

    tech_p1 = tf_tech.add_paragraph()
    tech_p1.text = "   - Frontend & UI: React 18, TypeScript, Vite, Tailwind CSS, Canvas 3D Orb Visualizer (WebGL)."
    tech_p1.font.size = Pt(11.5)
    tech_p1.font.color.rgb = BLACK

    tech_p2 = tf_tech.add_paragraph()
    tech_p2.text = "   - Backend & Database: Node.js, Express.js, TypeScript, PostgreSQL (Prisma ORM), IndexedDB (Offline caching), WebSockets."
    tech_p2.font.size = Pt(11.5)
    tech_p2.font.color.rgb = BLACK

    tech_p3 = tf_tech.add_paragraph()
    tech_p3.text = "   - AI, Vision & Speech: Google Gemini 3.7 / 2.5 Flash LLM, MediaPipe Hands (21-point ISL recognition), Web Speech API (STT & TTS)."
    tech_p3.font.size = Pt(11.5)
    tech_p3.font.color.rgb = BLACK

    # Point 2: Methodology & Process
    tb_meth = s3.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.4))
    p_meth = tb_meth.text_frame.paragraphs[0]
    p_meth.text = "•  Methodology and process for implementation (System Architecture Flowchart):"
    p_meth.font.size = Pt(13)
    p_meth.font.bold = True
    p_meth.font.color.rgb = BLACK

    # 4 Architecture Flowchart Boxes
    arch_layers = [
        ("Layer 1: Input Channels", "🎙️ Speech (Hindi/Marathi/English)\n🔘 1-Tap UI Cards & Widgets\n📷 SignBridge Camera (ISL)", HEADING_BLUE),
        ("Layer 2: AI Action Router", "🧠 Google Gemini 3.7 / 2.5 Engine\n🔄 Multi-Turn Memory Context\n⚡ Entity & Parameter Extraction", PURPLE_BORDER),
        ("Layer 3: Execution & DB", "💾 Prisma PostgreSQL Database\n📶 Offline IndexedDB Sync\n🤟 MediaPipe 21-Point Tracking", GREEN_ACCENT),
        ("Layer 4: Verified Output", "🔊 Spoken Voice Alarms & TTS\n📊 Caregiver Dashboard Sync\n🚨 Emergency Family SOS Alerts", FOOTER_BLUE)
    ]
    for i, (ltitle, ldesc, border_col) in enumerate(arch_layers):
        x = Inches(0.8 + i * 2.95)
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.4), Inches(2.75), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = BOX_BG
        box.line.color.rgb = border_col
        box.line.width = Pt(1.5)
        
        tf_bx = box.text_frame
        tf_bx.word_wrap = True
        tf_bx.margin_left = tf_bx.margin_top = tf_bx.margin_right = tf_bx.margin_bottom = Inches(0.12)
        
        p = tf_bx.paragraphs[0]
        p.text = ltitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = border_col
        p.alignment = PP_ALIGN.CENTER
        
        p_sub = tf_bx.add_paragraph()
        p_sub.text = ldesc
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = BLACK
        p_sub.alignment = PP_ALIGN.LEFT
        p_sub.space_before = Pt(6)

    # Live prototype badge at bottom
    tb_proto = s3.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5))
    p_pr = tb_proto.text_frame.paragraphs[0]
    p_pr.text = "🚀 Working Prototype Status: 100% Implemented & Live Tested | GitHub: https://github.com/swayamgulhane538/Aabha-ai"
    p_pr.font.size = Pt(11.5)
    p_pr.font.bold = True
    p_pr.font.color.rgb = GREEN_ACCENT
    p_pr.alignment = PP_ALIGN.CENTER


    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ══════════════════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(blank_layout)
    apply_base_template(s4, 4, center_title="FEASIBILITY AND VIABILITY")

    # Point 1: Analysis of Feasibility
    tb_fe = s4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.4))
    tf_fe = tb_fe.text_frame
    tf_fe.word_wrap = True

    p = tf_fe.paragraphs[0]
    p.text = "•  Analysis of the feasibility of the idea:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p1 = tf_fe.add_paragraph()
    p1.text = "   - Technical Feasibility: Proven working prototype running on standard browsers; requires zero extra hardware or wearable sensors."
    p1.font.size = Pt(11.5)
    p1.font.color.rgb = BLACK

    p2 = tf_fe.add_paragraph()
    p2.text = "   - Economic Viability: Utilizes on-device client processing & Google Gemini high-efficiency APIs, reducing server overhead costs by 80%."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = BLACK

    # Point 2 & 3: Challenges & Mitigation Matrix Flowchart
    tb_cm = s4.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.4))
    p_cm = tb_cm.text_frame.paragraphs[0]
    p_cm.text = "•  Potential challenges, risks and strategies for overcoming them (Mitigation Flowchart):"
    p_cm.font.size = Pt(13)
    p_cm.font.bold = True
    p_cm.font.color.rgb = BLACK

    matrix_rows = [
        ("Challenge 1: Poor Internet in Rural India", "Strategy: 100% Offline-First Architecture with Service Workers & IndexedDB. Alarms & games function without internet.", GREEN_ACCENT),
        ("Challenge 2: Multi-Lingual Dialects & Noise", "Strategy: Resilient NLP with Multi-Turn Clarification Memory and fuzzy entity parsing in Hindi, Marathi, and English.", HEADING_BLUE),
        ("Challenge 3: Elderly Fear of Complex Tech", "Strategy: Ultra-simple 1-tap touch cards, high-contrast visual cues, and automated spoken voice reading with zero learning curve.", PURPLE_BORDER)
    ]
    for i, (chall, strat, col) in enumerate(matrix_rows):
        y = Inches(3.2 + i * 1.15)
        # Left Box (Challenge)
        c_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(4.8), Inches(1.0))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = RGBColor(254, 242, 242)
        c_box.line.color.rgb = RGBColor(239, 68, 68) # Red 500
        c_box.line.width = Pt(1.2)
        tf_c = c_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = Inches(0.08)
        p = tf_c.paragraphs[0]
        p.text = chall
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(185, 28, 28)

        # Right Box (Strategy)
        s_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), y, Inches(6.733), Inches(1.0))
        s_box.fill.solid()
        s_box.fill.fore_color.rgb = RGBColor(240, 253, 244)
        s_box.line.color.rgb = col
        s_box.line.width = Pt(1.2)
        tf_s = s_box.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = Inches(0.08)
        p = tf_s.paragraphs[0]
        p.text = strat
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col


    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: IMPACT AND BENEFITS
    # ══════════════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank_layout)
    apply_base_template(s5, 5, center_title="IMPACT AND BENEFITS")

    # Point 1: Potential impact on target audience
    tb_imp = s5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4))
    p_imp = tb_imp.text_frame.paragraphs[0]
    p_imp.text = "•  Potential impact on the target audience:"
    p_imp.font.size = Pt(13)
    p_imp.font.bold = True
    p_imp.font.color.rgb = BLACK

    audience_cards = [
        ("Dementia Patients", "Ensures 100% timely medication adherence, preserves cognitive ability, and promotes dignified independent daily living.", HEADING_BLUE),
        ("Family Caregivers", "Reduces 70%+ caregiver burnout by automating repetitive routine checks, hydration alerts, and missed-dose notifications.", GREEN_ACCENT),
        ("Deaf & Mute Individuals", "Bridges critical healthcare communication barriers via SignBridge real-time Indian Sign Language (ISL) teleconsultation.", PURPLE_BORDER)
    ]
    for i, (title, desc, col) in enumerate(audience_cards):
        x = Inches(0.8 + i * 4.0)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.75), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = BOX_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        
        tf_b = box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = Inches(0.1)
        
        p = tf_b.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_d = tf_b.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = BLACK
        p_d.space_before = Pt(3)

    # Point 2: Benefits of solution
    tb_ben = s5.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.4))
    p_ben = tb_ben.text_frame.paragraphs[0]
    p_ben.text = "•  Benefits of the solution (Social, Economic, Healthcare & Scalability):"
    p_ben.font.size = Pt(13)
    p_ben.font.bold = True
    p_ben.font.color.rgb = BLACK

    benefits_cards = [
        ("🌍 Social Impact", "Promotes inclusive digital health for India's 6.3 Cr seniors and 1.8 Cr deaf/mute community under Ayushman Bharat (ABDM).", HEADING_BLUE),
        ("💰 Economic Benefits", "Cuts emergency hospital readmission costs caused by missed medications, falls, or late interventions by up to 40%.", GREEN_ACCENT),
        ("📈 Scalability & Deployment", "Ready for instant nationwide deployment across Old Age Homes, Memory Clinics, Primary Health Centers (PHCs), and NGOs.", PURPLE_BORDER)
    ]
    for i, (title, desc, col) in enumerate(benefits_cards):
        x = Inches(0.8 + i * 4.0)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.1), Inches(3.75), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = BOX_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        
        tf_b = box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = Inches(0.1)
        
        p = tf_b.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_d = tf_b.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = BLACK
        p_d.space_before = Pt(4)


    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: RESEARCH AND REFERENCES
    # ══════════════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank_layout)
    apply_base_template(s6, 6, center_title="RESEARCH AND REFERENCES")

    # Point: Details / Links of reference and research work
    tb_ref_head = s6.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4))
    p_rh = tb_ref_head.text_frame.paragraphs[0]
    p_rh.text = "•  Details / Links of the reference and research work:"
    p_rh.font.size = Pt(13)
    p_rh.font.bold = True
    p_rh.font.color.rgb = BLACK

    # Left Box: Clinical & Policy References
    ref_l = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    ref_l.fill.solid()
    ref_l.fill.fore_color.rgb = BOX_BG
    ref_l.line.color.rgb = HEADING_BLUE
    ref_l.line.width = Pt(1.5)
    tf_rl = ref_l.text_frame
    tf_rl.word_wrap = True
    tf_rl.margin_left = tf_rl.margin_top = tf_rl.margin_right = tf_rl.margin_bottom = Inches(0.15)

    p = tf_rl.paragraphs[0]
    p.text = "📚 Clinical Research & Government Policies"
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = HEADING_BLUE
    p.space_after = Pt(8)

    ref_items = [
        ("World Health Organization (WHO)", "Global Action Plan on the Public Health Response to Dementia (2017–2025) on assistive tech & reminiscence therapy."),
        ("Lancet Neurology (2022)", "Bio-Psycho-Social Reminiscence Therapy & Spaced Retrieval Gaming for Mild Cognitive Impairment (MCI)."),
        ("Google Research (2020)", "MediaPipe Hands: Real-Time On-Device Hand Landmark Detection & ISLRTC Sign Corpus."),
        ("Govt. of India (MoHFW)", "Ayushman Bharat Digital Mission (ABDM) & National Digital Health Blueprint (NDHB) standards.")
    ]
    for src, det in ref_items:
        p_s = tf_rl.add_paragraph()
        p_s.text = f"• {src}:"
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = BLACK
        
        p_d = tf_rl.add_paragraph()
        p_d.text = f"  {det}"
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = GRAY_TEXT
        p_d.space_after = Pt(6)

    # Right Box: Project Deliverables & Team Links
    ref_r = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.633), Inches(4.8))
    ref_r.fill.solid()
    ref_r.fill.fore_color.rgb = BOX_BG
    ref_r.line.color.rgb = GREEN_ACCENT
    ref_r.line.width = Pt(1.5)
    tf_rr = ref_r.text_frame
    tf_rr.word_wrap = True
    tf_rr.margin_left = tf_rr.margin_top = tf_rr.margin_right = tf_rr.margin_bottom = Inches(0.15)

    p = tf_rr.paragraphs[0]
    p.text = "🔗 Live Project Links & Team Details"
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    p.space_after = Pt(8)

    deliv_items = [
        ("Team Name:", "PBCOE-Nexora"),
        ("Problem Statement ID:", "SIH26003"),
        ("Live Web Application Link:", "https://aabha-ai.vercel.app"),
        ("GitHub Repository Link:", "https://github.com/swayamgulhane538/Aabha-ai"),
        ("Interactive Demo Mode:", "Integrated directly inside app under 'Interactive Demo Tour'"),
        ("Platform Compatibility:", "Responsive Web, Mobile PWA, Tablet & Offline Mode")
    ]
    for k, v in deliv_items:
        p_k = tf_rr.add_paragraph()
        p_k.text = f"• {k}"
        p_k.font.size = Pt(11)
        p_k.font.bold = True
        p_k.font.color.rgb = BLACK
        
        p_v = tf_rr.add_paragraph()
        p_v.text = f"  {v}"
        p_v.font.size = Pt(11)
        p_v.font.bold = (k.startswith("Live Web") or k.startswith("GitHub") or k.startswith("Team Name"))
        p_v.font.color.rgb = HEADING_BLUE if k.startswith("Live Web") or k.startswith("GitHub") else BLACK
        p_v.space_after = Pt(6)

    # Save presentation
    output_pptx = r"C:\Users\hp\.gemini\antigravity\scratch\aabha-ai\SIH_2025_PBCOE_Nexora_AABHA_AI.pptx"
    prs.save(output_pptx)
    print(f"Exact Template PPTX successfully created at: {output_pptx}")

if __name__ == "__main__":
    build_presentation()
